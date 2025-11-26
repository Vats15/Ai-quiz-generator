# app.py
import os
import streamlit as st

# ---------------------------
# Ensure OPENAI_API_KEY is available to modules imported later.
# Copy it from Streamlit secrets into the environment (if present).
# This avoids import-time failures in generator.py when secrets aren't yet in env.
# ---------------------------
if "OPENAI_API_KEY" not in os.environ:
    try:
        env_key = st.secrets.get("OPENAI_API_KEY")
        if env_key:
            os.environ["OPENAI_API_KEY"] = env_key
    except Exception:
        # If st.secrets isn't available for any reason, continue.
        pass

# Now safe to import generator which uses os.getenv("OPENAI_API_KEY")
from generator import generate_questions_from_text
from exporter import questions_to_dataframe, df_to_csv_bytes, questions_to_json_bytes

# Standard libs for file handling
from io import BytesIO

st.set_page_config(page_title="AI Quiz Generator", layout="wide")


# ---------------------------
# File parsing helpers (lazy import heavy libs)
# ---------------------------
def extract_text_from_pdf(uploaded_file) -> str:
    """
    Extract text from uploaded PDF file. Uses pdfplumber if available.
    Falls back to an empty string and error message if pdfplumber not installed.
    """
    uploaded_file.seek(0)
    try:
        import pdfplumber
    except Exception:
        st.error("PDF parsing requires `pdfplumber` installed. Ask the maintainer to add it.")
        return ""

    text_parts = []
    try:
        with pdfplumber.open(uploaded_file) as pdf:
            for p in pdf.pages:
                t = p.extract_text()
                if t:
                    text_parts.append(t)
    except Exception as e:
        st.error(f"Error parsing PDF: {e}")
        return ""
    return "\n".join(text_parts)


def extract_text_from_pptx(uploaded_file) -> str:
    """
    Extract text from uploaded PPTX file using python-pptx.
    If the package is missing, show error and return empty string.
    """
    uploaded_file.seek(0)
    try:
        from pptx import Presentation
    except Exception:
        st.error("PPTX parsing requires `python-pptx` installed. Ask the maintainer to add it.")
        return ""

    parts = []
    try:
        prs = Presentation(uploaded_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
    except Exception as e:
        st.error(f"Error parsing PPTX: {e}")
        return ""
    return "\n".join(parts)


# ---------------------------
# UI
# ---------------------------
st.title("AI Quiz Generator — MVP")

with st.sidebar:
    st.header("Input & Options")
    input_mode = st.radio("Input type", ["Paste text", "Upload PDF", "Upload PPTX"])
    if input_mode == "Paste text":
        source_text = st.text_area("Paste source text (or paste slides/notes)", height=300)
    else:
        file_type = "pdf" if input_mode == "Upload PDF" else "pptx"
        uploaded = st.file_uploader("Upload file", type=file_type)
        source_text = ""
        if uploaded is not None:
            if input_mode == "Upload PDF":
                source_text = extract_text_from_pdf(uploaded)
            else:
                source_text = extract_text_from_pptx(uploaded)

    qtype = st.selectbox("Question type", ["mcq", "tf", "full", "mixed"])
    n_questions = st.slider("Number of questions", 1, 30, 5)
    difficulty = st.selectbox("Difficulty", ["easy", "medium", "hard", "auto"])
    use_timer = st.checkbox("Enable timer (minutes)")
    timer_minutes = int(st.number_input("Timer duration (minutes)", min_value=1, max_value=180, value=10)) if use_timer else None

    st.markdown("---")
    generate_btn = st.button("Generate Questions")

# Initialize session state
if "questions" not in st.session_state:
    st.session_state.questions = []
if "timer_running" not in st.session_state:
    st.session_state.timer_running = False
if "timer_end" not in st.session_state:
    st.session_state.timer_end = None

# ---------------------------
# Generate button action
# ---------------------------
if generate_btn:
    if not source_text or len(source_text.strip()) < 40:
        st.error("Please provide source text (paste or upload a file) with enough content.")
    else:
        with st.spinner("Generating questions..."):
            try:
                qs = generate_questions_from_text(source_text, qtype, n_questions, difficulty)
                # normalize ids
                for i, q in enumerate(qs, start=1):
                    q.setdefault("id", i)
                st.session_state.questions = qs
                st.success(f"Generated {len(qs)} questions.")
                # start timer if requested
                if use_timer:
                    import time
                    st.session_state.timer_end = time.time() + (timer_minutes * 60)
                    st.session_state.timer_running = True
            except Exception as e:
                st.error(f"Generation failed: {e}")

# Timer display & logic (sidebar)
if use_timer and st.session_state.get("timer_running") and st.session_state.get("timer_end"):
    import time
    remaining = int(st.session_state.timer_end - time.time())
    if remaining <= 0:
        st.session_state.timer_running = False
        st.success("Time's up! You can now reveal answers.")
    else:
        mins, secs = divmod(remaining, 60)
        st.sidebar.info(f"Time remaining: {mins:02d}:{secs:02d}")
        # lightweight refresh to update timer
        st.experimental_rerun()

# ---------------------------
# Main UI layout
# ---------------------------
cols = st.columns((3, 1))
with cols[0]:
    st.header("Questions")
    if not st.session_state.questions:
        st.info("No questions yet — enter source text and click Generate.")
    else:
        for q in st.session_state.questions:
            qid = q.get("id")
            qtype_local = q.get("type")
            excerpt = q.get("question", "")[:120].replace("\n", " ")
            with st.expander(f"{qid}. {excerpt}...", expanded=False):
                st.write(q.get("question"))
                if qtype_local == "mcq":
                    opts = q.get("options", [])
                    if opts:
                        for idx, o in enumerate(opts):
                            st.write(f"{chr(65+idx)}. {o}")
                if not use_timer or (use_timer and not st.session_state.get("timer_running")):
                    if st.button(f"Show Answer for Q{qid}", key=f"show_{qid}"):
                        st.markdown(f"**Answer:** {q.get('answer')}")
                        if q.get("explanation"):
                            st.markdown(f"**Explanation:** {q.get('explanation')}")
                else:
                    st.info("Answers are hidden while the timer is running.")

with cols[1]:
    st.header("Actions")
    if st.session_state.questions:
        df = questions_to_dataframe(st.session_state.questions)
        csv_bytes = df_to_csv_bytes(df)
        json_bytes = questions_to_json_bytes(st.session_state.questions)
        st.download_button("Download CSV", data=csv_bytes, file_name="questions.csv", mime="text/csv")
        st.download_button("Download JSON", data=json_bytes, file_name="questions.json", mime="application/json")
        st.markdown("---")
        st.write("Preview")
        st.dataframe(df)
        if st.button("Clear questions"):
            st.session_state.questions = []
            st.experimental_rerun()
    else:
        st.write("No actions available yet.")
